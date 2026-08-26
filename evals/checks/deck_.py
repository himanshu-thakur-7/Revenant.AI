"""Pitch-deck (.pptx) checks — structural, via python-pptx directly on the
file, matching the contract agents/sales/prompt.py's SALES_SYSTEM actually
states (title/bullet length limits) rather than inventing a new one.
"""

from __future__ import annotations

import re
from pathlib import Path

from evals.checks import Check

_PLACEHOLDER_MARKERS = ("Lorem ipsum", "TODO", "{{", "<company>", "<merchant>", "[COMPANY]")
_CTA_VERBS = ("book", "schedule", "start", "try", "get", "let's", "reply", "connect", "talk")


def _open(pptx_path: str):
    if not pptx_path or not Path(pptx_path).exists():
        return None
    try:
        from pptx import Presentation
        return Presentation(pptx_path)
    except Exception:  # noqa: BLE001
        return None


def pptx_opens(pptx_path: str, *, name: str = "pptx_opens") -> Check:
    if not pptx_path or not Path(pptx_path).exists():
        return Check(name, False, f"no file at {pptx_path or '(empty path)'}")
    prs = _open(pptx_path)
    return Check(name, prs is not None, "opened cleanly" if prs else "python-pptx raised on open")


def _slide_texts(prs) -> list[list[str]]:
    out = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs).strip()
                    if t:
                        texts.append(t)
        out.append(texts)
    return out


def slide_count_between(pptx_path: str, lo: int = 5, hi: int = 7,
                        *, name: str = "slide_count") -> Check:
    prs = _open(pptx_path)
    if prs is None:
        return Check(name, False, "deck did not open")
    n = len(prs.slides)
    return Check(name, lo <= n <= hi, f"{n} slides (want {lo}-{hi})", measured=n)


def slide_arc(pptx_path: str, startup: str, merchant: str, *, name: str = "slide_arc") -> Check:
    prs = _open(pptx_path)
    if prs is None:
        return Check(name, False, "deck did not open")
    slides = _slide_texts(prs)
    if not slides:
        return Check(name, False, "no slides")
    first = " ".join(slides[0]).lower()
    last = " ".join(slides[-1]).lower()
    names_present = (startup or "").lower() in first and (merchant or "").lower() in first
    cta_present = any(v in last for v in _CTA_VERBS)
    ok = names_present and cta_present
    detail = f"slide1 names both companies: {names_present}; last slide has a CTA verb: {cta_present}"
    return Check(name, ok, detail)


def copy_limits(pptx_path: str, max_title_words: int = 8, max_bullet_words: int = 15,
                *, name: str = "copy_limits") -> Check:
    prs = _open(pptx_path)
    if prs is None:
        return Check(name, False, "deck did not open")
    violations = []
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            is_title = shape == slide.shapes.title if slide.shapes.title else False
            limit = max_title_words if is_title else max_bullet_words
            for para in shape.text_frame.paragraphs:
                t = "".join(run.text for run in para.runs).strip()
                if not t:
                    continue
                n = len(t.split())
                if n > limit:
                    violations.append(f"slide{i}:{'title' if is_title else 'bullet'} "
                                     f"{n}w > {limit}w: {t[:50]!r}")
    return Check(name, not violations, "; ".join(violations[:5]) if violations else "within limits",
                 measured=len(violations))


def no_placeholder_text(pptx_path: str, *, name: str = "no_placeholder") -> Check:
    prs = _open(pptx_path)
    if prs is None:
        return Check(name, False, "deck did not open")
    all_text = " ".join(" ".join(t) for t in _slide_texts(prs))
    hits = [m for m in _PLACEHOLDER_MARKERS if m.lower() in all_text.lower()]
    return Check(name, not hits, f"found: {hits}" if hits else "clean", measured=hits)
