"""T0 coverage for evals/checks/deck_.py::slide_arc — pure python-pptx
logic against real files this test builds itself, no network. Regression
guard for a real false negative found this session: PhonePe's genuine,
on-brief closing-ask slide ("THE ASK" / "Next Steps:" / "Explore ...",
"See ...", "Enhance ...") failed slide_arc because none of those verbs
were in the original 9-word _CTA_VERBS list — Meesho's near-identical
slide only passed by the luck of also containing "Let's" and "Schedule".
See evals/checks/deck_.py's own comment on the fix.
"""

from __future__ import annotations

from pptx import Presentation

from evals.checks.deck_ import slide_arc


def _deck(tmp_path, first_texts, last_texts, n_middle=4):
    prs = Presentation()
    layout = prs.slide_layouts[6]  # blank layout — full control over text boxes

    def add_slide(texts):
        slide = prs.slides.add_slide(layout)
        box = slide.shapes.add_textbox(0, 0, 9144000, 1000000)
        tf = box.text_frame
        tf.text = texts[0]
        for t in texts[1:]:
            p = tf.add_paragraph()
            p.text = t

    add_slide(first_texts)
    for i in range(n_middle):
        add_slide([f"Middle slide {i}"])
    add_slide(last_texts)

    p = tmp_path / "deck.pptx"
    prs.save(str(p))
    return str(p)


def test_phonepe_style_ask_slide_now_passes(tmp_path):
    # The actual real-world failure this test guards against regressing.
    path = _deck(
        tmp_path,
        first_texts=["Razorpay x PhonePe"],
        last_texts=["THE ASK", "Next Steps: Explore Razorpay for PhonePe",
                    "See live prototype", "Explore integration possibilities",
                    "Enhance financial operations"],
    )
    c = slide_arc(path, "Razorpay", "PhonePe")
    assert c.passed, c.detail


def test_meesho_style_ask_slide_passes(tmp_path):
    path = _deck(
        tmp_path,
        first_texts=["Razorpay x Meesho"],
        last_texts=["THE ASK", "Let's Transform Your Payment Experience",
                    "Explore a live prototype", "See results on your data",
                    "Schedule a quick call"],
    )
    c = slide_arc(path, "Razorpay", "Meesho")
    assert c.passed, c.detail


def test_missing_both_company_names_fails(tmp_path):
    path = _deck(tmp_path, first_texts=["A generic pitch deck"],
                last_texts=["Thanks! Let's talk soon."])
    c = slide_arc(path, "Razorpay", "Meesho")
    assert not c.passed


def test_no_real_ask_on_last_slide_fails(tmp_path):
    # No CTA verb anywhere on a genuinely content-free closing slide.
    path = _deck(tmp_path, first_texts=["Razorpay x Meesho"],
                last_texts=["Thank you.", "Questions?"])
    c = slide_arc(path, "Razorpay", "Meesho")
    assert not c.passed
