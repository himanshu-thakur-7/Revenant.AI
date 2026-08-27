"""evals/checks/deck_.py — the leave-behind deck.

Builds real .pptx files with python-pptx rather than mocking, because
the checks read shape/paragraph structure and a mock would let a real
structural regression through.
"""

from __future__ import annotations

import pytest

pptx = pytest.importorskip("pptx")

import evals.checks.deck_ as deck_
from pptx import Presentation


def build(tmp_path, slides, name="deck.pptx"):
    """slides: list of list-of-strings, one list per slide."""
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for texts in slides:
        s = prs.slides.add_slide(blank)
        box = s.shapes.add_textbox(0, 0, 9144000, 1200000)
        tf = box.text_frame
        tf.text = texts[0] if texts else ""
        for t in texts[1:]:
            tf.add_paragraph().text = t
    p = tmp_path / name
    prs.save(str(p))
    return str(p)


def good_deck(tmp_path):
    return build(tmp_path, [
        ["Razorpay x PhonePe"],
        ["The problem", "Reconciliation is manual", "Errors compound"],
        ["Why now", "UPI volume is growing"],
        ["How it works", "Parse settlement files", "Match the ledger"],
        ["Proof", "Live prototype on your data"],
        ["THE ASK", "Next Steps", "See the live prototype", "Schedule a call"],
    ])


# ── opening ───────────────────────────────────────────────────────────

def test_a_real_deck_opens(tmp_path):
    assert deck_.pptx_opens(good_deck(tmp_path)).passed


def test_a_missing_file_fails(tmp_path):
    c = deck_.pptx_opens(str(tmp_path / "nope.pptx"))
    assert not c.passed and "no file" in c.detail


def test_an_empty_path_fails():
    assert not deck_.pptx_opens("").passed


def test_a_corrupt_pptx_fails_cleanly(tmp_path):
    # A truncated upload/mux produces exactly this: right name, not a deck.
    p = tmp_path / "broken.pptx"
    p.write_bytes(b"not a zip archive at all")
    c = deck_.pptx_opens(str(p))
    assert not c.passed
    assert "raised on open" in c.detail


def test_downstream_checks_fail_closed_on_a_corrupt_deck(tmp_path):
    p = tmp_path / "broken.pptx"
    p.write_bytes(b"garbage")
    for fn in (deck_.slide_count_between, deck_.copy_limits, deck_.no_placeholder_text):
        assert not fn(str(p)).passed


# ── slide count ───────────────────────────────────────────────────────

def test_a_six_slide_deck_passes(tmp_path):
    assert deck_.slide_count_between(good_deck(tmp_path)).passed


def test_a_too_short_deck_fails(tmp_path):
    # Two slides means the Sales step gave up partway.
    c = deck_.slide_count_between(build(tmp_path, [["A"], ["B"]]))
    assert not c.passed


def test_a_too_long_deck_fails(tmp_path):
    c = deck_.slide_count_between(build(tmp_path, [[f"S{i}"] for i in range(20)]))
    assert not c.passed


def test_slide_count_boundaries_are_inclusive(tmp_path):
    assert deck_.slide_count_between(build(tmp_path, [[f"S{i}"] for i in range(5)])).passed
    assert deck_.slide_count_between(build(tmp_path, [[f"S{i}"] for i in range(7)])).passed


# ── copy limits ───────────────────────────────────────────────────────

def test_concise_copy_passes(tmp_path):
    assert deck_.copy_limits(good_deck(tmp_path)).passed


def test_an_essay_bullet_fails(tmp_path):
    # A 40-word bullet is unreadable on a slide and signals the model
    # wrote prose instead of a deck.
    long_bullet = " ".join(["word"] * 40)
    c = deck_.copy_limits(build(tmp_path, [["Title"], ["Heading", long_bullet]]))
    assert not c.passed
    assert c.measured >= 1


def test_the_violation_detail_names_the_slide(tmp_path):
    c = deck_.copy_limits(build(tmp_path, [["T"], ["H", " ".join(["w"] * 40)]]))
    assert "slide2" in c.detail


def test_only_the_first_few_violations_are_reported(tmp_path):
    # The detail string goes into a prompt and a chat message; an
    # unbounded list would swamp both.
    many = [["T"]] + [["H", " ".join(["w"] * 40)] for _ in range(10)]
    assert len(deck_.copy_limits(build(tmp_path, many)).detail) < 1200


def test_empty_paragraphs_are_ignored(tmp_path):
    assert deck_.copy_limits(build(tmp_path, [["Title", "", "   ", "ok bullet"]])).passed


# ── placeholder text ──────────────────────────────────────────────────

@pytest.mark.parametrize("ph", ["Lorem ipsum dolor", "TODO: fill this in",
                                "{{company}}", "<company> value prop", "[COMPANY]"])
def test_placeholder_text_is_caught(tmp_path, ph):
    c = deck_.no_placeholder_text(build(tmp_path, [["Title"], ["Heading", ph]]))
    assert not c.passed
    assert c.measured


def test_a_clean_deck_has_no_placeholders(tmp_path):
    assert deck_.no_placeholder_text(good_deck(tmp_path)).passed


def test_placeholder_detection_is_case_insensitive(tmp_path):
    assert not deck_.no_placeholder_text(build(tmp_path, [["T"], ["lorem IPSUM"]])).passed


# ── slide arc ─────────────────────────────────────────────────────────

def test_a_real_arc_passes(tmp_path):
    assert deck_.slide_arc(good_deck(tmp_path), "Razorpay", "PhonePe").passed


def test_slide_one_must_name_both_companies(tmp_path):
    d = build(tmp_path, [["A generic pitch"], ["x"], ["x"], ["x"], ["x"],
                         ["THE ASK", "Schedule a call"]])
    assert not deck_.slide_arc(d, "Razorpay", "PhonePe").passed


def test_the_last_slide_must_actually_ask(tmp_path):
    d = build(tmp_path, [["Razorpay x PhonePe"], ["x"], ["x"], ["x"], ["x"],
                         ["Thank you.", "Questions?"]])
    assert not deck_.slide_arc(d, "Razorpay", "PhonePe").passed


def test_arc_on_an_empty_deck_fails(tmp_path):
    prs = Presentation()
    p = tmp_path / "empty.pptx"
    prs.save(str(p))
    c = deck_.slide_arc(str(p), "Razorpay", "PhonePe")
    assert not c.passed
    assert "no slides" in c.detail
