"""Pre-built, co-branded **Razorpay × boAt** pitch deck for the on-stage demo.

Deterministic — generates a polished 5-slide .pptx in the Razorpay palette with
a Razorpay × boAt co-brand lockup on every slide. If real logo PNGs are dropped
at ``agents/demo_razorpay_assets/logos/{razorpay,boat}.png`` they're used;
otherwise clean brand-styled wordmarks are drawn instead.

Regenerate:  python -m agents.demo_razorpay_deck
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Razorpay palette ──────────────────────────────────────────────
BLUE = RGBColor(0x33, 0x95, 0xFF)
NAVY = RGBColor(0x0D, 0x23, 0x66)
INK = RGBColor(0x0D, 0x23, 0x66)
MUTED = RGBColor(0x5A, 0x6B, 0x87)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF1, 0xF6, 0xFF)
GREEN = RGBColor(0x14, 0x9E, 0x5A)
CORAL = RGBColor(0xE1, 0x25, 0x3B)  # boAt accent

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

_ASSETS = Path(__file__).parent / "demo_razorpay_assets"
_LOGO_RZP = _ASSETS / "logos" / "razorpay.png"
_LOGO_BOAT = _ASSETS / "logos" / "boat.png"
DECK_PATH = _ASSETS / "razorpay-boat-deck.pptx"


def _bg(slide, color=WHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _rect(slide, x, y, w, h, color, *, shape=MSO_SHAPE.RECTANGLE, line=None):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def _text(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
          space_after=6):
    """runs: list of (text, size, color, bold) tuples → one paragraph each."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        if isinstance(para, list):  # multiple runs in one paragraph
            for (t, size, color, bold) in para:
                r = p.add_run(); r.text = t
                r.font.size = Pt(size); r.font.color.rgb = color
                r.font.bold = bold; r.font.name = "Calibri"
        else:
            t, size, color, bold = para
            r = p.add_run(); r.text = t
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = "Calibri"
    return tb


def _cobrand(slide, *, y=Inches(0.32), small=False):
    """Razorpay × boAt lockup — real PNGs if present, else styled wordmarks."""
    x = Inches(0.55)
    h = Inches(0.42) if not small else Inches(0.34)
    # Razorpay
    if _LOGO_RZP.exists():
        slide.shapes.add_picture(str(_LOGO_RZP), x, y, height=h)
        rx = x + Inches(1.7)
    else:
        _rect(slide, x, y, Inches(0.34), Inches(0.34), BLUE,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        _text(slide, x - Inches(0.02), y - Inches(0.02), Inches(0.4), Inches(0.4),
              [("R", 15, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, x + Inches(0.42), y - Inches(0.03), Inches(1.5), Inches(0.42),
              [("Razorpay", 17, NAVY, True)], anchor=MSO_ANCHOR.MIDDLE)
        rx = x + Inches(1.85)
    # ×
    _text(slide, rx, y - Inches(0.03), Inches(0.35), Inches(0.42),
          [("×", 16, MUTED, False)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    bx = rx + Inches(0.35)
    # boAt
    if _LOGO_BOAT.exists():
        slide.shapes.add_picture(str(_LOGO_BOAT), bx, y, height=h)
    else:
        _text(slide, bx, y - Inches(0.03), Inches(1.4), Inches(0.42),
              [[("bo", 17, INK, True), ("A", 17, CORAL, True), ("t", 17, INK, True)]],
              anchor=MSO_ANCHOR.MIDDLE)


def _footer(slide, idx, total):
    _text(slide, Inches(11.4), Inches(7.02), Inches(1.6), Inches(0.35),
          [(f"Razorpay × boAt   ·   {idx:02d}/{total:02d}", 9, MUTED, False)],
          align=PP_ALIGN.RIGHT)


# ── slides ────────────────────────────────────────────────────────
def _title_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, NAVY)
    _rect(s, 0, Inches(6.9), SLIDE_W, Inches(0.6), BLUE)
    # co-brand (light on navy)
    x, y = Inches(0.55), Inches(0.45)
    _rect(s, x, y, Inches(0.4), Inches(0.4), BLUE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _text(s, x - Inches(0.01), y - Inches(0.02), Inches(0.42), Inches(0.44),
          [("R", 17, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _text(s, x + Inches(0.5), y - Inches(0.03), Inches(2), Inches(0.46),
          [("Razorpay", 19, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
    _text(s, x + Inches(2.1), y - Inches(0.03), Inches(0.4), Inches(0.46),
          [("×", 17, RGBColor(0x8F, 0xA6, 0xD8), False)], anchor=MSO_ANCHOR.MIDDLE)
    _text(s, x + Inches(2.5), y - Inches(0.03), Inches(2), Inches(0.46),
          [[("bo", 19, WHITE, True), ("A", 19, RGBColor(0xFF, 0x8A, 0x9A), True),
            ("t", 19, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    _text(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(0.5),
          [("A WORKING PROTOTYPE + PILOT FOR boAt", 15, BLUE, True)])
    _text(s, Inches(0.86), Inches(3.0), Inches(11.6), Inches(2.2),
          [("Turn boAt's checkout drop-off", 44, WHITE, True),
           ("into recovered revenue.", 44, RGBColor(0x8F, 0xC2, 0xFF), True)],
          space_after=2)
    _text(s, Inches(0.9), Inches(5.2), Inches(10.5), Inches(1),
          [("Razorpay Magic Checkout — 1-click, address-prefilled, RTO-protected "
            "checkout for the 100M+ shoppers already saved with Razorpay.", 16,
            RGBColor(0xC7, 0xD6, 0xF5), False)])


def _leak_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _cobrand(s); _footer(s, 2, 5)
    _text(s, Inches(0.9), Inches(1.15), Inches(11), Inches(0.4),
          [("THE LEAK", 13, BLUE, True)])
    _text(s, Inches(0.86), Inches(1.55), Inches(11.6), Inches(1),
          [("boAt loses revenue at the last step", 34, NAVY, True)])
    cards = [
        ("Cart drop-off", "Manual address + payment entry on mobile is the #1 "
         "reason boAt shoppers abandon a full cart."),
        ("COD return-to-origin", "Cash-on-delivery on high-value audio drives "
         "costly RTO — packed, shipped, refused, returned."),
        ("After-hours support", "“Where is my order?” load piles up when "
         "no one's at the desk to answer."),
    ]
    cw = Inches(3.75); gap = Inches(0.28); x0 = Inches(0.9); y0 = Inches(2.9)
    for i, (h, b) in enumerate(cards):
        x = x0 + i * (cw + gap)
        _rect(s, x, y0, cw, Inches(3.1), LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        _rect(s, x, y0, cw, Inches(0.12), CORAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        _text(s, x + Inches(0.3), y0 + Inches(0.4), cw - Inches(0.6), Inches(0.6),
              [(h, 19, NAVY, True)])
        _text(s, x + Inches(0.3), y0 + Inches(1.15), cw - Inches(0.6), Inches(1.8),
              [(b, 14, MUTED, False)])


def _solution_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _cobrand(s); _footer(s, 3, 5)
    _text(s, Inches(0.9), Inches(1.15), Inches(11), Inches(0.4),
          [("THE FIX — MAGIC CHECKOUT", 13, BLUE, True)])
    _text(s, Inches(0.86), Inches(1.55), Inches(11.6), Inches(1),
          [("One tap to buy. Zero drop-off.", 34, NAVY, True)])
    rows = [
        ("1-click prefilled checkout",
         "100M+ shoppers are already saved with Razorpay — on boAt they're "
         "recognised and the checkout fills itself. One tap to pay."),
        ("Predictive COD & RTO Protection",
         "Risk models nudge shaky COD orders to prepaid and reimburse failed "
         "deliveries — boAt keeps more of every high-value order."),
        ("100+ payment methods + offers",
         "UPI, cards, netbanking, wallets, EMI, BNPL, plus a coupon engine — the "
         "method the shopper is most likely to complete, auto-recommended."),
        ("Live in a day",
         "One script tag on boAt's store — Shopify, WooCommerce, or custom. "
         "PCI-DSS Level 1. No re-platforming."),
    ]
    y = Inches(2.75)
    for i, (h, b) in enumerate(rows):
        yy = y + i * Inches(1.02)
        _rect(s, Inches(0.9), yy + Inches(0.05), Inches(0.14), Inches(0.72), BLUE,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        _text(s, Inches(1.25), yy, Inches(4.4), Inches(0.9),
              [(h, 18, NAVY, True)], anchor=MSO_ANCHOR.MIDDLE)
        _text(s, Inches(5.9), yy, Inches(6.6), Inches(0.9),
              [(b, 13.5, MUTED, False)], anchor=MSO_ANCHOR.MIDDLE)


def _impact_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, NAVY)
    _text(s, Inches(0.9), Inches(0.9), Inches(11), Inches(0.4),
          [("THE IMPACT FOR boAt", 13, RGBColor(0x8F, 0xC2, 0xFF), True)])
    _text(s, Inches(0.86), Inches(1.3), Inches(11.6), Inches(1),
          [("What one tap recovers", 34, WHITE, True)])
    stats = [("+40%", "conversion uplift", "less form-filling drop-off"),
             ("100M+", "pre-saved shoppers", "recognised on boAt instantly"),
             ("↓ RTO", "return-to-origin", "predictive COD → prepaid"),
             ("~1 day", "to go live", "one script tag, no re-platform")]
    cw = Inches(2.85); gap = Inches(0.26); x0 = Inches(0.9); y0 = Inches(2.9)
    for i, (big, lab, sub) in enumerate(stats):
        x = x0 + i * (cw + gap)
        _rect(s, x, y0, cw, Inches(2.9), RGBColor(0x14, 0x2E, 0x6E),
              shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        _text(s, x, y0 + Inches(0.5), cw, Inches(1),
              [(big, 40, RGBColor(0x8F, 0xC2, 0xFF), True)], align=PP_ALIGN.CENTER)
        _text(s, x, y0 + Inches(1.6), cw, Inches(0.5),
              [(lab, 15, WHITE, True)], align=PP_ALIGN.CENTER)
        _text(s, x + Inches(0.2), y0 + Inches(2.05), cw - Inches(0.4), Inches(0.7),
              [(sub, 11.5, RGBColor(0xB8, 0xCB, 0xF0), False)], align=PP_ALIGN.CENTER)


def _cta_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _cobrand(s); _footer(s, 5, 5)
    _rect(s, 0, Inches(6.9), SLIDE_W, Inches(0.6), BLUE)
    _text(s, Inches(0.9), Inches(2.2), Inches(11), Inches(0.4),
          [("THE ASK", 13, BLUE, True)])
    _text(s, Inches(0.86), Inches(2.65), Inches(11.6), Inches(1.4),
          [("A 30-day Magic Checkout pilot", 40, NAVY, True),
           ("on one boAt storefront — live this week.", 40, BLUE, True)],
          space_after=2)
    _text(s, Inches(0.9), Inches(4.7), Inches(11), Inches(0.8),
          [[("We'll ship the integration, wire COD/RTO protection, and report "
             "recovered revenue against a baseline. ", 16, MUTED, False),
            ("You only scale what pays.", 16, NAVY, True)]])
    _rect(s, Inches(0.9), Inches(5.7), Inches(3.3), Inches(0.7), BLUE,
          shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _text(s, Inches(0.9), Inches(5.7), Inches(3.3), Inches(0.7),
          [("Book the pilot →", 17, WHITE, True)], align=PP_ALIGN.CENTER,
          anchor=MSO_ANCHOR.MIDDLE)


def build_deck(out_path: Path | None = None) -> Path:
    out = Path(out_path) if out_path else DECK_PATH
    prs = Presentation()
    prs.slide_width = Emu(int(SLIDE_W)); prs.slide_height = Emu(int(SLIDE_H))
    _title_slide(prs); _leak_slide(prs); _solution_slide(prs)
    _impact_slide(prs); _cta_slide(prs)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


if __name__ == "__main__":
    p = build_deck()
    print("deck written:", p, p.stat().st_size, "bytes")
