"""Render a pitch deck (`.pptx`) from a list of slides.

The LLM composes the outline; this module handles the mechanical layout.
Design is stripped down and on-brand for a technical B2B pitch — dark
background, teal (`#52e0c4`) accents on titles and rules, JetBrains-Mono
eyebrow labels, Space-Grotesk body text (falls back to sans-serif when
the reader doesn't have the fonts).
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt


# Shroud/Revenant palette (matches the console theme + prototypes).
BG = RGBColor(0x05, 0x06, 0x0A)
PANEL = RGBColor(0x0C, 0x0F, 0x17)
INK = RGBColor(0xE6, 0xEB, 0xF5)
MUTED = RGBColor(0x8A, 0x94, 0xA8)
ACCENT = RGBColor(0x52, 0xE0, 0xC4)  # wisp teal
AMBER = RGBColor(0xF5, 0xB9, 0x42)


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def render_deck(slides: list[dict[str, Any]], out_path: Path,
                *, deck_title: str) -> Path:
    """Render slides into a `.pptx` at ``out_path``.

    Each slide dict may contain:
        title   (str)          — required
        bullets (list[str])    — optional
        notes   (str)          — optional, becomes speaker notes
        kind    (str)          — 'title' | 'content' (default) | 'cta'
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for i, s in enumerate(slides):
        kind = s.get("kind") or ("title" if i == 0 else "content")
        title = s.get("title", "")
        bullets = s.get("bullets", []) or []
        notes = s.get("notes", "") or ""

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        _paint_background(slide, BG)
        _accent_bar(slide)

        if kind == "title":
            _title_slide(slide, title=title, subtitle=deck_title,
                         eyebrow=s.get("eyebrow", "REVENANT"))
        elif kind == "cta":
            _cta_slide(slide, title=title, bullets=bullets)
        else:
            _content_slide(slide, title=title, bullets=bullets,
                           eyebrow=s.get("eyebrow", f"{i+1:02d} / {len(slides):02d}"))

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


# ── layout helpers ────────────────────────────────────────────
def _paint_background(slide, color: RGBColor) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False


def _accent_bar(slide) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.6), Inches(0.85),
                                 Inches(0.06), Inches(0.32))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def _text_box(slide, x, y, w, h, text: str, *, size: int, color: RGBColor,
              bold: bool = False, font: str = "Space Grotesk"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return box


def _title_slide(slide, *, title: str, subtitle: str, eyebrow: str) -> None:
    _text_box(slide, Inches(0.85), Inches(0.8), Inches(11), Inches(0.4),
              eyebrow.upper(), size=12, color=ACCENT, bold=True,
              font="JetBrains Mono")

    tf = slide.shapes.add_textbox(Inches(0.85), Inches(2.4),
                                  Inches(11.5), Inches(2.6)).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(60); r.font.bold = True
    r.font.color.rgb = INK; r.font.name = "Space Grotesk"

    _text_box(slide, Inches(0.85), Inches(5.8), Inches(11.5), Inches(0.5),
              subtitle, size=20, color=MUTED)


def _content_slide(slide, *, title: str, bullets: list[str],
                   eyebrow: str) -> None:
    _text_box(slide, Inches(0.85), Inches(0.85), Inches(6), Inches(0.4),
              eyebrow.upper(), size=11, color=ACCENT, bold=True,
              font="JetBrains Mono")

    tf = slide.shapes.add_textbox(Inches(0.85), Inches(1.4),
                                  Inches(11.5), Inches(1.4)).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(40); r.font.bold = True
    r.font.color.rgb = INK; r.font.name = "Space Grotesk"

    # bullets
    body = slide.shapes.add_textbox(Inches(0.85), Inches(3.15),
                                    Inches(11.5), Inches(3.9))
    tf = body.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i else tf.paragraphs[0]
        r = p.add_run()
        r.text = f"◆  {b}"
        r.font.size = Pt(22)
        r.font.color.rgb = INK if i == 0 else RGBColor(0xCB, 0xD1, 0xE0)
        r.font.name = "Space Grotesk"
        p.space_after = Pt(14)


def _cta_slide(slide, *, title: str, bullets: list[str]) -> None:
    _text_box(slide, Inches(0.85), Inches(0.85), Inches(6), Inches(0.4),
              "THE ASK", size=11, color=AMBER, bold=True,
              font="JetBrains Mono")

    tf = slide.shapes.add_textbox(Inches(0.85), Inches(1.8),
                                  Inches(11.5), Inches(2.2)).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(52); r.font.bold = True
    r.font.color.rgb = INK; r.font.name = "Space Grotesk"

    body = slide.shapes.add_textbox(Inches(0.85), Inches(4.6),
                                    Inches(11.5), Inches(2.4))
    tf = body.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i else tf.paragraphs[0]
        r = p.add_run(); r.text = f"→  {b}"
        r.font.size = Pt(22); r.font.color.rgb = ACCENT if i == 0 else INK
        r.font.name = "Space Grotesk"; r.font.bold = i == 0
        p.space_after = Pt(12)
